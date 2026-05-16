"""Extract plain text from PDF, DOCX, PPTX, XLSX, CSV, and plain text."""
import csv
import io


def extract_text(content: bytes, ext: str) -> str:
    """Dispatch to the right parser based on file extension."""
    handlers = {
        ".pdf":  _from_pdf,
        ".docx": _from_docx,
        ".pptx": _from_pptx,
        ".xlsx": _from_xlsx,
        ".csv":  _from_csv,
        ".txt":  _from_txt,
    }
    return handlers.get(ext.lower(), _from_txt)(content)


def _from_pdf(content: bytes) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(stream=content, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)


def _from_docx(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _from_pptx(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def _from_xlsx(content: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    rows = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            line = ", ".join(str(c) for c in row if c is not None)
            if line:
                rows.append(line)
    return "\n".join(rows)


def _from_csv(content: bytes) -> str:
    reader = csv.reader(io.StringIO(content.decode("utf-8", errors="replace")))
    return "\n".join(", ".join(row) for row in reader)


def _from_txt(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")
